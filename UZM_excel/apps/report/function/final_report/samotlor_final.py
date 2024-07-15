import os
import io
import matplotlib.pyplot as plt
import numpy as np

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from Field.models import Run, Well
from report.function.model_service import get_data, get_waste
from excel_parcer.models import Data

from report.models import Graf2Param, Graf4Param, Graf3Param, Graf6Param, Graf1Param
import matplotlib
matplotlib.use('agg')

def SamotlorTitle(prs, my_well):
    """Заполняет информацией титульный слайд"""
    title_slide = prs.slides[0]
    title_table = title_slide.shapes[4].table

    cell = title_table.cell(0, 1)
    cell.text = str(my_well.pad_name.field.field_name)
    cell = title_table.cell(1, 1)
    cell.text = str(my_well.pad_name.pad_name)
    cell = title_table.cell(2, 1)
    cell.text = str(my_well.well_name)
    cell = title_table.cell(3, 1)
    cell.text = str(my_well.get_well_type())
    cell = title_table.cell(4, 1)
    cell.text = '-' if my_well.T3_end is None else my_well.T3_end.strftime('%d-%m-%Y')

    # Форматируем текст всех ячеек
    for index in range(5):
        my_cell = title_table.cell(index, 1)
        for p in my_cell.text_frame.paragraphs:
            p.text = p.text.upper()
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.LEFT


def SamotlorInfo(prs, well_obj, my_wellbore):
    """Заполняет информацией слайд основной информации"""
    info_slide = prs.slides[1]
    info_table = info_slide.shapes[1].table

    cell = info_table.cell(0, 1)
    cell.text = str(well_obj.pad_name.field.client)
    cell = info_table.cell(1, 1)
    cell.text = str(well_obj.pad_name.field)
    cell = info_table.cell(2, 1)
    cell.text = well_obj.pad_name.pad_name
    cell = info_table.cell(3, 1)
    cell.text = well_obj.well_name
    cell = info_table.cell(4, 1)
    cell.text = '-' if well_obj.coordinate_system is None else well_obj.coordinate_system
    cell = info_table.cell(5, 1)
    cell.text = '-' if well_obj.latitude is None else well_obj.latitude
    cell = info_table.cell(6, 1)
    cell.text = '-' if well_obj.longtitude is None else well_obj.longtitude
    cell = info_table.cell(7, 1)
    cell.text = '-' if well_obj.NY is None else str(well_obj.NY)
    cell = info_table.cell(8, 1)
    cell.text = '-' if well_obj.EX is None else str(well_obj.EX)
    cell = info_table.cell(9, 1)
    cell.text = '-' if well_obj.RKB is None else str(well_obj.RKB)
    cell = info_table.cell(10, 1)
    cell.text = my_wellbore.get_contractor_name()
    cell = info_table.cell(0, 3)
    cell.text = '-' if well_obj.geomagnetic_model is None else well_obj.geomagnetic_model
    cell = info_table.cell(1, 3)
    cell.text = '-' if well_obj.geomagnetic_date is None else well_obj.geomagnetic_date.strftime('%d-%m-%Y')
    cell = info_table.cell(2, 3)
    cell.text = '-' if well_obj.btotal is None else str(well_obj.btotal)
    cell = info_table.cell(3, 3)
    cell.text = '-' if well_obj.dip is None else str(well_obj.dip)
    cell = info_table.cell(4, 3)
    cell.text = '-' if well_obj.dec is None else str(well_obj.dec)
    cell = info_table.cell(5, 3)
    cell.text = '-' if well_obj.grid_convergence is None else str(well_obj.grid_convergence)
    cell = info_table.cell(6, 3)
    cell.text = '-' if well_obj.gtotal is None else str(well_obj.gtotal)
    cell = info_table.cell(7, 3)
    cell.text = well_obj.get_north_direction()
    cell = info_table.cell(8, 3)
    if well_obj.T1_end is not None and well_obj.T1_start is not None:
        cell.text = well_obj.T1_start.strftime('%d.%m.%Y') + '-' + well_obj.T1_end.strftime('%d.%m.%Y')
    else:
        cell.text = '-'
    cell = info_table.cell(9, 3)
    if well_obj.T3_end is not None and well_obj.T3_start is not None:
        cell.text = well_obj.T3_start.strftime('%d.%m.%Y') + '-' + well_obj.T3_end.strftime('%d.%m.%Y')
    else:
        cell.text = '-'
    cell = info_table.cell(10, 3)
    cell.text = '-' if well_obj.сontractor is None else well_obj.сontractor.drill_contractor_name

    # Форматируем текст всех ячеек
    for my_cell in info_table.iter_cells():
        for p in my_cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.CENTER


def SamotlorQualityCharts(prs, run_obj):
    """Заполняет информацией слайдов контроля качества"""
    # !!! Данные для графика !!!
    wellbore_obj = run_obj.section.wellbore
    runs = Run.objects.filter(section__wellbore=wellbore_obj)

    context = {'depth': list(),
               'depthGtotal': list(),
               'depthGref': list(),
               'depthGmax': list(),
               'depthGmin': list(),
               'depthDipraw': list(),
               'depthDipref': list(),
               'depthDipmax': list(),
               'depthDipmin': list(),
               'depthDipcorr': list(),
               "depthBoxy": list(),
               "depthBz": list(),
               "depthGoxy": list(),
               "depthGz": list(),
               "depthBtotal": list(),
               "depthBref": list(),
               "depthBmax": list(),
               "depthBmin": list(),
               "depthBcorr": list(),
               }
    for run in runs:
        surveys = Data.objects.filter(run=run).order_by('depth')
        for survey in surveys:
            context['depth'].append(survey.depth)
            # График Gtotal
            context['depthGtotal'].append(survey.Gtotal())
            context['depthGref'].append(wellbore_obj.well_name.gtotal)
            context['depthGmax'].append(wellbore_obj.well_name.max_gtotal())
            context['depthGmin'].append(wellbore_obj.well_name.min_gtotal())
            # График Boxy-Bz
            context['depthBoxy'].append(survey.get_boxy())
            context['depthBz'].append(survey.BZ)
            # График Btotal
            context['depthBtotal'].append(survey.Btotal())
            context['depthBref'].append(wellbore_obj.well_name.btotal)
            context['depthBmax'].append(wellbore_obj.well_name.max_btotal())
            context['depthBmin'].append(wellbore_obj.well_name.min_btotal())
            context['depthBcorr'].append(survey.Btotal_corr)
            # График Dip
            context['depthDipraw'].append(survey.Dip())
            context['depthDipref'].append(wellbore_obj.well_name.dip)
            context['depthDipmax'].append(wellbore_obj.well_name.max_dip())
            context['depthDipmin'].append(wellbore_obj.well_name.min_dip())
            context['depthDipcorr'].append(survey.DIP_corr)

    # отступы графика от краёв
    margins = {  # +++
        "left": 0.09,
        "bottom": 0.2,
        "right": 0.85,
        "top": 0.90
    }
    bbox = (1.002, 1)

    # отступы легенды графика от угла
    # !!! Графики контроля качества 1/2 !!!
    graph_slide = prs.slides[2]
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(8.73, 2.6))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthGtotal'], color='#D6D64C', label='Gt')
    plt.plot(context['depth'], context['depthGref'], color='g', label='Gt_REF')
    plt.plot(context['depth'], context['depthGmax'], '--r')
    plt.plot(context['depth'], context['depthGmin'], '--r')
    plt.title('График напряженности гравитационного поля', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Gtotal, G', color='gray', fontsize=10)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf2Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf2Param.DoesNotExist:
        pass
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(0.6), Inches(0.8)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)

    # График магнитного наклонения
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(8.73, 2.6))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthDipraw'], color='#D6D64C', label='Dip_RAW')
    plt.plot(context['depth'], context['depthDipref'], color='green', label='Dip_REF')
    plt.plot(context['depth'], context['depthDipmax'], '--r')
    plt.plot(context['depth'], context['depthDipmin'], '--r')
    plt.plot(context['depth'], context['depthDipcorr'], color='blue', label='Dip_CORR')
    plt.title('График магнитного наклонения', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Магнитное наклонение, гр', color='gray', fontsize=10)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf6Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf6Param.DoesNotExist:
        pass
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(0.6), Inches(3.6)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)

    # !!! Графики контроля качества 2/2 !!!
    graph_slide = prs.slides[3]
    # График соотношения показаний осевого и поперечных магнитометров
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(8.73, 2.6))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    # print('x=', context['depth'])
    # print('y=', context['depthBz'])

    plt.plot(context['depth'], context['depthBoxy'], color='blue', label='Boxy')
    plt.plot(context['depth'], context['depthBz'], color='red', label='Bz')
    plt.title('График соотношения показаний осевого и поперечного магнитометров', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Boxy, нТл', color='gray', fontsize=10)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf1Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf1Param.DoesNotExist:
        pass
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(0.6), Inches(0.8)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)

    # График напряженности геомагнитного поля
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(8.73, 2.6))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthBtotal'], color='#D6D64C', label='Bt_RAW')
    plt.plot(context['depth'], context['depthBcorr'], color='blue', label='Bt_CORR')
    plt.plot(context['depth'], context['depthBmax'], '--r')
    plt.plot(context['depth'], context['depthBref'], color='green', label='Bt_REF')
    plt.plot(context['depth'], context['depthBmin'], '--r')
    plt.title('График напряженности геомагнитного поля', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Напряженность геомагнитного поля, нТл', color='gray', fontsize=6)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf4Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf4Param.DoesNotExist:
        pass
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(0.6), Inches(3.6)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)


def SamotlorProj(prs, run_obj):
    """Заполняет информацией слайдов проекции"""
    projection_slide = prs.slides[4]
    wellbore_obj = run_obj.section.wellbore
    img_path = os.getcwd() + f"\\Files\\Report_out\\{wellbore_obj}.png"
    left = Inches(0.2)
    top = Inches(0.9)
    pic = projection_slide.shapes.add_picture(img_path, left, top, width=Inches(9.3), height=Inches(6.2))
    pic.crop_bottom = 0.30
    pic.crop_left = 0.07
    pic.crop_right = 0.07
    pic.crop_top = 0.08


def SamotlorRecomendation(prs, my_wellbore):
    """ Слайд с выводами и рекомендациями"""
    waste = get_waste(my_wellbore)
    slide = prs.slides[7]
    slide.shapes[0].text_frame.paragraphs[2].text = (slide.shapes[0].text_frame.paragraphs[2].text
                                                     .replace('{{meas_dot}}', str(waste[0]['Точка замера']))
                                                     .replace('{{hor_waste}}', str(waste[0]['Отход по горизонтали']))
                                                     .replace('{{hor_word}}', str(waste[1]['hor'])))
    slide.shapes[0].text_frame.paragraphs[3].text = (slide.shapes[0].text_frame.paragraphs[3].text
                                                     .replace('{{ver_waste}}', str(abs(waste[0]['Отход по вертикали'])))
                                                     .replace('{{ver_word}}', str(waste[1]['ver'])))
    slide.shapes[0].text_frame.paragraphs[4].text = (slide.shapes[0].text_frame.paragraphs[4].text
                                                     .replace('{{comm_waste}}', str(waste[0]['Общий отход'])))

    # Форматируем текст
    slide.shapes[0].text_frame.paragraphs[2].font.size = Pt(12)
    slide.shapes[0].text_frame.paragraphs[2].font.name = 'Segoe UI'
    slide.shapes[0].text_frame.paragraphs[3].font.size = Pt(12)
    slide.shapes[0].text_frame.paragraphs[3].font.name = 'Segoe UI'
    slide.shapes[0].text_frame.paragraphs[4].font.size = Pt(12)
    slide.shapes[0].text_frame.paragraphs[4].font.name = 'Segoe UI'


