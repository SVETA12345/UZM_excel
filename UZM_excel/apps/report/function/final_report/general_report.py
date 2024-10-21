import os
import io

import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from Field.models import Run
from report.function.model_service import get_waste
from excel_parcer.models import Data

from .samotlor_final import SamotlorTitle, SamotlorInfo, SamotlorRecomendation, SamotlorQualityCharts, SamotlorProj
from report.models import Graf2Param, Graf4Param, Graf6Param, Graf3Param
import matplotlib
matplotlib.use('agg')

def FRCreator(run_id: int) -> tuple[str, str] | str:
    """ Конструктор финальных отчётов"""
    # Объекты моделей
    my_run = Run.objects.get(id=run_id)
    param = {"my_run": my_run,
             "my_well": my_run.section.wellbore.well_name,
             "my_wellbore": my_run.section.wellbore}
    if my_run.section.wellbore.well_name.pad_name.field.client.client_name == 'SMTL':
        return samotlor_final(**param)
    return standard_final(**param)


def standard_final(my_run, my_well, my_wellbore) -> str:
    """ Конструктор стандартных отчётов """
    # Шаблон всех итоговых отчётов
    if my_run.section.wellbore.well_name.pad_name.field.client.client_name == 'UGNG':
        prs = Presentation(os.getcwd() + "\\files\\Шаблон\\Юганск_итоговый.pptx")
    else:
        prs = Presentation(os.getcwd() + "\\files\\Шаблон\\Итоговый_отчет.pptx")
    # !!! Титульный слайд !!!
    TitleMaker(prs, my_well)
    # !!! Общая информация !!!
    InfoMaker(prs, my_well, my_wellbore)
    # !!! Графики контроля качества !!!
    QualityChartsMaker(prs, my_run)
    # !!! Рекомендации !!!
    RecomendationMaker(prs, my_wellbore)
    # !!! Проекции !!!
    ProjMaker(prs, my_run)

    file_dir = os.getcwd() + "\\files\\Report_out\\final_report.pptx"
    prs.save(file_dir)
    return file_dir


def samotlor_final(my_run, my_well, my_wellbore) -> str:
    """ Конструктор отчётов для Самотлора """
    if my_well.get_well_type() == 'ВНС':
        prs = Presentation(os.getcwd() + "\\files\\Шаблон\\Самотлорское_итоговый_ВНС.pptx")
    else:
        prs = Presentation(os.getcwd() + "\\files\\Шаблон\\Самотлорское_итоговый.pptx")
    # !!! Титульный слайд !!!
    SamotlorTitle(prs, my_well)
    # !!! Общая информация !!!
    SamotlorInfo(prs, my_well, my_wellbore)
    # !!! Графики контроля качества !!!
    SamotlorQualityCharts(prs, my_run)
    # !!! Рекомендации !!!
    SamotlorRecomendation(prs, my_wellbore)
    # !!! Проекции !!!
    SamotlorProj(prs, my_run)

    file_dir = os.getcwd() + "\\files\\Report_out\\final_report.pptx"
    prs.save(file_dir)
    return file_dir


def TitleMaker(prs, my_well):
    """Заполняет информацией титульный слайд"""
    title_slide = prs.slides[1]
    title_slide.shapes[1].text_frame.add_paragraph()
    p_well = title_slide.shapes[1].text_frame.add_paragraph()
    if my_well.pad_name.field.field_name.find('кое') != -1:
        p_well.text = f'по скважине {my_well.well_name} куста {my_well.pad_name.pad_name} ' + \
                      f'{my_well.pad_name.field.field_name} месторождения'.replace('кое', 'кого')
    else:
        p_well.text = f'по скважине {my_well.well_name} куста {my_well.pad_name.pad_name} ' + \
                      f'месторождения {my_well.pad_name.field.field_name}'

    title_slide.shapes[1].text_frame.add_paragraph()
    p_info = title_slide.shapes[1].text_frame.add_paragraph()
    p_info.text = 'Независимый контроль качества телеметрических замеров в процессе бурения'

    # Форматируем вставленный текст
    for p in title_slide.shapes[1].text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.name = 'Segoe UI'
        p.alignment = PP_ALIGN.CENTER


def InfoMaker(prs, well_obj, my_wellbore):
    """Заполняет информацией слайд основной информации"""
    info_slide = prs.slides[2]
    info_table = info_slide.shapes[4].table

    cell = info_table.cell(0, 1)
    cell.text = str(well_obj.pad_name.field.client)
    cell = info_table.cell(1, 1)
    cell.text = str(well_obj.pad_name.field)
    cell = info_table.cell(2, 1)
    cell.text = well_obj.pad_name.pad_name
    cell = info_table.cell(3, 1)
    cell.text = well_obj.well_name
    cell = info_table.cell(4, 1)
    cell.text = '-' if well_obj.coordinate_system is None else well_obj.coordinate_system
    cell = info_table.cell(5, 1)
    cell.text = '-' if well_obj.latitude is None else well_obj.latitude
    cell = info_table.cell(6, 1)
    cell.text = '-' if well_obj.longtitude is None else well_obj.longtitude
    cell = info_table.cell(7, 1)
    cell.text = '-' if well_obj.NY is None else str(well_obj.NY)
    cell = info_table.cell(8, 1)
    cell.text = '-' if well_obj.EX is None else str(well_obj.EX)
    cell = info_table.cell(9, 1)
    cell.text = '-' if well_obj.RKB is None else str(well_obj.RKB)
    cell = info_table.cell(10, 1)
    cell.text = my_wellbore.get_contractor_name()
    cell = info_table.cell(0, 3)
    cell.text = '-' if well_obj.geomagnetic_model is None else well_obj.geomagnetic_model
    cell = info_table.cell(1, 3)
    cell.text = '-' if well_obj.geomagnetic_date is None else well_obj.geomagnetic_date.strftime('%d-%m-%Y')
    cell = info_table.cell(2, 3)
    cell.text = '-' if well_obj.btotal is None else str(well_obj.btotal)
    cell = info_table.cell(3, 3)
    cell.text = '-' if well_obj.dip is None else str(well_obj.dip)
    cell = info_table.cell(4, 3)
    cell.text = '-' if well_obj.dec is None else str(well_obj.dec)
    cell = info_table.cell(5, 3)
    cell.text = '-' if well_obj.grid_convergence is None else str(well_obj.grid_convergence)
    cell = info_table.cell(6, 3)
    cell.text = '-' if well_obj.gtotal is None else str(well_obj.gtotal)
    cell = info_table.cell(7, 3)
    cell.text = well_obj.get_north_direction()
    cell = info_table.cell(8, 3)
    cell.text = '-'
    cell = info_table.cell(9, 3)
    cell.text = '-'

    # Форматируем текст всех ячеек
    for my_cell in info_table.iter_cells():
        for p in my_cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.CENTER


def QualityChartsMaker(prs, run_obj):
    """Заполняет информацией слайдов контроля качества"""
    # !!! Данные для графика !!!
    wellbore_obj = run_obj.section.wellbore
    runs = Run.objects.filter(section__wellbore=wellbore_obj)

    context = {'depth': list(),
               'depthGtotal': list(),
               'depthGref': list(),
               'depthGmax': list(),
               'depthGmin': list(),
               'depthDipraw': list(),
               'depthDipref': list(),
               'depthDipmax': list(),
               'depthDipmin': list(),
               'depthDipcorr': list(),
               "depthBoxy": list(),
               "depthBz": list(),
               "depthGoxy": list(),
               "depthGz": list(),
               "depthBtotal": list(),
               "depthBref": list(),
               "depthBmax": list(),
               "depthBmin": list(),
               "depthBcorr": list(),
               }
    for run in runs:
        surveys = Data.objects.filter(run=run).order_by('depth')
        for survey in surveys:
            context['depth'].append(survey.depth)
            # График Gtotal
            context['depthGtotal'].append(survey.Gtotal())
            context['depthGref'].append(wellbore_obj.well_name.gtotal)
            context['depthGmax'].append(wellbore_obj.well_name.max_gtotal())
            context['depthGmin'].append(wellbore_obj.well_name.min_gtotal())
            # График Goxy-Gz
            context['depthGoxy'].append(survey.get_goxy())
            context['depthGz'].append(survey.CZ)
            # График Boxy-Bz
            context['depthBoxy'].append(survey.get_boxy())
            context['depthBz'].append(survey.BZ)
            # График Btotal
            context['depthBtotal'].append(survey.Btotal())
            context['depthBref'].append(wellbore_obj.well_name.btotal)
            context['depthBmax'].append(wellbore_obj.well_name.max_btotal())
            context['depthBmin'].append(wellbore_obj.well_name.min_btotal())
            context['depthBcorr'].append(survey.Btotal_corr)
            # График Dip
            context['depthDipraw'].append(survey.Dip())
            context['depthDipref'].append(wellbore_obj.well_name.dip)
            context['depthDipmax'].append(wellbore_obj.well_name.max_dip())
            context['depthDipmin'].append(wellbore_obj.well_name.min_dip())
            context['depthDipcorr'].append(survey.DIP_corr)

    # отступы графика от краёв
    margins = {  # +++
        "left": 0.09,
        "bottom": 0.2,
        "right": 0.85,
        "top": 0.90
    }
    # отступы легенды графика от угла
    bbox = (1.02, 1)

    # !!! Графики контроля качества 1/2 !!!
    graph_slide = prs.slides[3]
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(11.1, 2.8))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthGtotal'], color='#D6D64C', label='Gtotal')
    plt.plot(context['depth'], context['depthGref'], color='g', label='Gref')
    plt.plot(context['depth'], context['depthGmax'], '--r')
    plt.plot(context['depth'], context['depthGmin'], '--r')
    plt.title('График напряженности гравитационного поля', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Gtotal, G', color='gray', fontsize=10)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf2Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf2Param.DoesNotExist:
        print('Данные графика не заполнены')
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(1.1), Inches(0.8)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)

    # График магнитного наклонения
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(11.1, 2.8))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthDipraw'], color='#D6D64C', label='DipRaw')
    plt.plot(context['depth'], context['depthDipref'], color='green', label='DipRef')
    plt.plot(context['depth'], context['depthDipmax'], '--r')
    plt.plot(context['depth'], context['depthDipmin'], '--r')
    plt.plot(context['depth'], context['depthDipcorr'], color='blue', label='DipCorr')
    plt.title('График магнитного наклонения', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Магнитное наклонение, гр', color='gray', fontsize=10)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf6Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf6Param.DoesNotExist:
        print('Данные графика не заполнены')
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(1.1), Inches(3.6)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)

    # !!! Графики контроля качества 2/2 !!!
    graph_slide = prs.slides[4]

    # График соотношения показаний осевого и поперечных магнитометров
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(1, 1, figsize=(11.1, 2.8))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthBoxy'], color='blue', label='Boxy')
    plt.plot(context['depth'], context['depthBz'], color='red', label='Bz')
    plt.title('График соотношения показаний осевого и поперечного магнитометров', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Boxy, нТл', color='gray', fontsize=10)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf3Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf3Param.DoesNotExist:
        print('Данные графика не заполнены')
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(1.1), Inches(0.8)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)

    # График напряженности геомагнитного поля
    image_stream = io.BytesIO()
    fig, ax = plt.subplots(figsize=(11.1, 2.8))
    currFig = id(fig)
    plt.subplots_adjust(**margins)
    plt.plot(context['depth'], context['depthBtotal'], color='#D6D64C', label='Bt_RAW')
    plt.plot(context['depth'], context['depthBcorr'], color='blue', label='Bcorr')
    plt.plot(context['depth'], context['depthBmax'], '--r')
    plt.plot(context['depth'], context['depthBref'], color='green', label='Bref')
    plt.plot(context['depth'], context['depthBmin'], '--r')
    plt.title('График напряженности геомагнитного поля', fontsize=12)
    plt.xlabel('Измеренная глубина, м', color='gray', fontsize=10)
    plt.ylabel('Напряженность геомагнитного поля, нТл', color='gray', fontsize=6)
    plt.legend(bbox_to_anchor=bbox, loc='upper left', fontsize=10)
    try:
        graph_param = Graf4Param.objects.get(wellbore=wellbore_obj)
        ax.set_xlim(graph_param.x_min, graph_param.x_max)
        ax.set_xticks(np.arange(graph_param.x_min, graph_param.x_max, graph_param.x_del))
        ax.set_ylim(graph_param.y_min, graph_param.y_max)
        ax.set_yticks(np.arange(graph_param.y_min, graph_param.y_max, graph_param.y_del))
    except Graf4Param.DoesNotExist:
        print('Данные графика не заполнены')
    plt.grid(True)
    plt.savefig(image_stream)
    x, y = Inches(1.1), Inches(3.6)
    pic = graph_slide.shapes.add_picture(image_stream, x, y)
    plt.close(currFig)


def ProjMaker(prs, run_obj):
    """Заполняет информацией слайдов проекции"""
    projection_slide = prs.slides[5]
    wellbore_obj = run_obj.section.wellbore
    runs = Run.objects.filter(section__wellbore=wellbore_obj)
    # all_data = get_data(runs)
    # get_graphics(all_data, run_obj.section.wellbore)

    img_path = os.getcwd() + f"\\Files\\Report_out\\{wellbore_obj}.png"
    left = Inches(2)
    top = Inches(0.8)
    pic = projection_slide.shapes.add_picture(img_path, left, top, width=Inches(9.3), height=Inches(6.2))
    pic.crop_bottom = 0.33
    pic.crop_left = 0.07
    pic.crop_right = 0.07
    pic.crop_top = 0.08
    # img_path = 'F:\\Рабочий стол\\GIT\\UZM_excel\\files\\Шаблон\\321.bmp'
    # left = Inches(6.4)
    # top = Inches(1)
    # pic = projection_slide.shapes.add_picture(img_path, left, top, width=Inches(5.8), height=Inches(5.8))


def RecomendationMaker(prs, my_wellbore):
    """ Слайд с выводами и рекомендациями"""
    waste = get_waste(my_wellbore)
    slide = prs.slides[6]
    slide.shapes[5].text_frame.paragraphs[2].text = (slide.shapes[5].text_frame.paragraphs[2].text
                                                     .replace('{{meas_dot}}', str(waste[0]['Точка замера']))
                                                     .replace('{{hor_waste}}', str(waste[0]['Отход по горизонтали']))
                                                     .replace('{{hor_word}}', str(waste[1]['hor'])))
    slide.shapes[5].text_frame.paragraphs[3].text = (slide.shapes[5].text_frame.paragraphs[3].text
                                                     .replace('{{ver_waste}}', str(abs(waste[0]['Отход по вертикали'])))
                                                     .replace('{{ver_word}}', str(waste[1]['ver'])))
    slide.shapes[5].text_frame.paragraphs[4].text = (slide.shapes[5].text_frame.paragraphs[4].text
                                                     .replace('{{comm_waste}}', str(waste[0]['Общий отход'])))

    # Форматируем текст
    slide.shapes[5].text_frame.paragraphs[2].font.size = Pt(13.5)
    slide.shapes[5].text_frame.paragraphs[2].font.name = 'Segoe UI'
    slide.shapes[5].text_frame.paragraphs[3].font.size = Pt(13.5)
    slide.shapes[5].text_frame.paragraphs[3].font.name = 'Segoe UI'
    slide.shapes[5].text_frame.paragraphs[4].font.size = Pt(13.5)
    slide.shapes[5].text_frame.paragraphs[4].font.name = 'Segoe UI'

