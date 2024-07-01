from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation("F:\\Рабочий стол\\GIT\\UZM_excel\\files\\Шаблон\\Самотлорское_итоговый.pptx")

# title_slide = prs.slides[0]
# title_table = title_slide.shapes[4].table
# for i, my_cell in enumerate(title_table.iter_cells()):
#     print(i, my_cell.text)


slide = prs.slides[6]
for i, shape in enumerate(slide.shapes):
    if not shape.has_text_frame:
        print(i, shape)
    else:
        for k, p in enumerate(shape.text_frame.paragraphs):
            for run in p.runs:
                print('shape:', i,'paragraph:',k, run.text)

# !!! Титульный слайд !!!

# title_slide = prs.slides[1]
# title_slide.shapes[3].text_frame.add_paragraph()
# p_well = title_slide.shapes[3].text_frame.add_paragraph()
# p_well.text = 'по скважине 271 куста 2'
#
# p_field = title_slide.shapes[3].text_frame.add_paragraph()
# p_field.text = 'Немчиновского месторождения'
#
# title_slide.shapes[3].text_frame.add_paragraph()
# p_info = title_slide.shapes[3].text_frame.add_paragraph()
# p_info.text = 'Независимый контроль качества телеметрических замеров в процессе бурения'
#
# Форматируем вставленный текст
# for p in title_slide.shapes[3].text_frame.paragraphs:
#     p.font.bold = True
#     p.font.size = Pt(20)
#     p.font.name = 'Segoe UI'
#     p.alignment = PP_ALIGN.CENTER

# !!! Общая информация !!!

# info_slide = prs.slides[2]
# info_table = info_slide.shapes[6].table
#
# cell = info_table.cell(0, 1)
# cell.text = 'Компания!'
# cell = info_table.cell(1, 1)
# cell.text = 'Месторождение!'
# cell = info_table.cell(2, 1)
# cell.text = 'Куст!'
# cell = info_table.cell(3, 1)
# cell.text = 'Скважина!'
# cell = info_table.cell(4, 1)
# cell.text = 'Система координат!'
# cell = info_table.cell(5, 1)
# cell.text = 'Широта!'
# cell = info_table.cell(6, 1)
# cell.text = 'Долгота!'
# cell = info_table.cell(7, 1)
# cell.text = 'NY!'
# cell = info_table.cell(8, 1)
# cell.text = 'EX!'
# cell = info_table.cell(9, 1)
# cell.text = 'Альтитуда точки отсчёта!'
# cell = info_table.cell(10, 1)
# cell.text = 'Подрядчик по ННБ!'
# cell = info_table.cell(0, 3)
# cell.text = 'Геомагнитная привязка'
# cell = info_table.cell(1, 3)
# cell.text = 'Дата гп'
# cell = info_table.cell(2, 3)
# cell.text = 'Напряжение гп'
# cell = info_table.cell(3, 3)
# cell.text = 'Магнитное наклонение'
# cell = info_table.cell(4, 3)
# cell.text = 'Магнитное склонение'
# cell = info_table.cell(5, 3)
# cell.text = 'Угол сближения меридеан'
# cell = info_table.cell(6, 3)
# cell.text = 'Напряденность грав. поля'
# cell = info_table.cell(7, 3)
# cell.text = 'Принятое направление'
# cell = info_table.cell(8, 3)
# cell.text = 'Данные локал поля магнитных аномалий'
# cell = info_table.cell(9, 3)
# cell.text = 'Данные внешнего гп'
#
# # Форматируем текст всех ячеек
# for my_cell in info_table.iter_cells():
#     for p in my_cell.text_frame.paragraphs:
#         p.font.bold = True
#         p.font.size = Pt(13)
#         p.font.name = 'Segoe UI'
#         p.alignment = PP_ALIGN.CENTER


# !!! Графики контроля качества 1/2 !!!
graph_slide = prs.slides[6]

chart_data = ChartData()
chart_data.add_series('Bt_RAW',    (32.2, 28.4, 34.7))
chart_data.add_series('Bt_res',    (24.3, 30.6, 20.2))
chart_data.add_series('Bt_Cor', (20.4, 18.3, 26.2))
chart_data.categories = (1, 2, 3)
# chart_data.x_values = [1, 2, 3]
# chart_data.Y_values = [2, 4, 6]
x, y, cx, cy = Inches(1), Inches(1), Inches(6), Inches(4)
graph_slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)

# !!! Графики контроля качества 2/2 !!!

# !!! Проекции !!!
# projection_slide = prs.slides[5]
#
# img_path = 'F:\\Рабочий стол\\GIT\\UZM_excel\\files\\Шаблон\\123.bmp'
# left = Inches(0.2)
# top = Inches(1)
# pic = projection_slide.shapes.add_picture(img_path, left, top, width=Inches(5.8), height=Inches(5.8))
#
# img_path = 'F:\\Рабочий стол\\GIT\\UZM_excel\\files\\Шаблон\\321.bmp'
# left = Inches(6.4)
# top = Inches(1)
# pic = projection_slide.shapes.add_picture(img_path, left, top, width=Inches(5.8), height=Inches(5.8))


prs.save("F:\\Загрузки\\Итоговый_отчет.pptx")
